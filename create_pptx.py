#!/usr/bin/env python3
"""Generate the AirBase Ops hackathon pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD   = RGBColor(0x16, 0x20, 0x3A)   # card bg
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # blue-500
ACCENT2   = RGBColor(0x22, 0xD3, 0xEE)   # cyan-400
GREEN     = RGBColor(0x22, 0xC5, 0x5E)   # green
ORANGE    = RGBColor(0xF9, 0x73, 0x16)   # orange
RED       = RGBColor(0xEF, 0x44, 0x44)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
YELLOW    = RGBColor(0xFA, 0xCC, 0x15)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────
def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val=str(int(alpha * 1000)))
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=16, color=WHITE,
                  bold=False, align=PP_ALIGN.LEFT, font_name="Calibri", line_spacing=1.2, bullet=False):
    """lines can be list of str or list of (str, color, bold, size)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            txt, c, b = line[0], line[1], line[2]
            s = line[3] if len(line) > 3 else size
        else:
            txt, c, b, s = line, color, bold, size
        if bullet and txt:
            p.text = txt
        else:
            p.text = txt
        p.font.size = Pt(s)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.8), color=ACCENT):
    return add_rect(slide, left, top, width, height, color)


def slide_number(slide, num, total):
    add_text(slide, W - Inches(1.2), H - Inches(0.55), Inches(1), Inches(0.4),
             f"{num}/{total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
solid_bg(sl, BG_DARK)

# Accent stripe top
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# Title
add_text(sl, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "AIRBASE OPS", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         font_name="Calibri")

# Subtitle
add_text(sl, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
         "AI-Powered Air Base Operations Simulator", size=28, color=ACCENT2,
         align=PP_ALIGN.CENTER)

# Tagline
add_text(sl, Inches(1.5), Inches(4.0), Inches(10), Inches(0.7),
         "Digitizing Saab's logistics training game with real-time AI decision support",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)

# Bottom info bar
add_rect(sl, Inches(0), H - Inches(1.0), W, Inches(1.0), BG_CARD)
add_text(sl, Inches(1.5), H - Inches(0.85), Inches(10), Inches(0.6),
         "Saab Smart Air Base Hackathon  |  March 2026  |  Tekniska Museet, Stockholm",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(sl, 1, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — The Problem
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "THE PROBLEM", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Air superiority is decided on the ground", size=36, bold=True, color=WHITE)

add_text(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
         "Sweden's dispersed air base commanders train with paper cards, physical dice, and whiteboard schedules.",
         size=18, color=LIGHT)

# Pain points — left column
pain_items = [
    ("Requires 8-15 people in one room", "Training happens 2-3 times per year"),
    ("Manual dice rolls & paper tracking", "A single game-day takes 4-8 real hours"),
    ("No data capture whatsoever", "Decisions evaporate after the session"),
]

for i, (title, desc) in enumerate(pain_items):
    y = Inches(2.7) + Inches(i * 1.3)
    add_accent_bar(sl, Inches(0.8), y, height=Inches(0.9), color=RED)
    add_text(sl, Inches(1.1), y, Inches(5), Inches(0.4), title, size=17, bold=True, color=WHITE)
    add_text(sl, Inches(1.1), y + Inches(0.4), Inches(5), Inches(0.4), desc, size=14, color=MUTED)

# Pain points — right column
pain_items2 = [
    ("No 'what-if' capability", "Commanders can't explore alternatives mid-game"),
    ("No AI decision support", "Learning only from experienced mentors (bottleneck)"),
    ("Single-player impossible", "Individuals can't practice on their own"),
]

for i, (title, desc) in enumerate(pain_items2):
    y = Inches(2.7) + Inches(i * 1.3)
    add_accent_bar(sl, Inches(7.0), y, height=Inches(0.9), color=RED)
    add_text(sl, Inches(7.3), y, Inches(5), Inches(0.4), title, size=17, bold=True, color=WHITE)
    add_text(sl, Inches(7.3), y + Inches(0.4), Inches(5), Inches(0.4), desc, size=14, color=MUTED)

# Bottom quote
add_rect(sl, Inches(0.6), H - Inches(1.2), Inches(12), Inches(0.7), BG_CARD)
add_text(sl, Inches(0.9), H - Inches(1.1), Inches(11.4), Inches(0.5),
         "\"In the next conflict, the bottleneck won't be aircraft — it will be the ground crews and logistics commanders who turn them around.\"",
         size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

slide_number(sl, 2, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — The Users
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "THE USERS", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Who uses this and what do they need?", size=36, bold=True, color=WHITE)

# Three user cards
users = [
    ("Basbatchef", "Base Commander", "Translates Air Tasking Orders into\nground-level decisions. Manages the\nbig picture: which aircraft fly, when,\nwith what weapons.", ACCENT),
    ("Underhallsberedare", "Maintenance Planner", "Allocates aircraft to tasks, plans\nrepair schedules, manages exchange\nunits and spare parts to keep the\nfleet mission-capable.", ACCENT2),
    ("Klargoringstropp", "Preparation Crews", "Execute hands-on fueling, arming,\nand pre-flight checks. Report faults.\nWork 8-hour rotating shifts under\ntime pressure.", GREEN),
]

for i, (swe, eng, desc, accent_c) in enumerate(users):
    x = Inches(0.8) + Inches(i * 4.0)
    # Card background
    add_rect(sl, x, Inches(2.0), Inches(3.6), Inches(4.0), BG_CARD)
    # Accent top
    add_rect(sl, x, Inches(2.0), Inches(3.6), Inches(0.06), accent_c)
    # Swedish title
    add_text(sl, x + Inches(0.3), Inches(2.3), Inches(3.0), Inches(0.5),
             swe, size=20, bold=True, color=accent_c)
    # English title
    add_text(sl, x + Inches(0.3), Inches(2.85), Inches(3.0), Inches(0.4),
             eng, size=14, color=MUTED)
    # Description
    add_text(sl, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(2.4),
             desc, size=14, color=LIGHT)

# Goal line
add_rect(sl, Inches(0.6), H - Inches(1.2), Inches(12), Inches(0.7), BG_CARD)
add_text(sl, Inches(0.9), H - Inches(1.15), Inches(11.4), Inches(0.25),
         "THEIR SHARED GOAL", size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.9), H - Inches(0.85), Inches(11.4), Inches(0.4),
         "Maximize sortie generation — get mission-capable Gripen aircraft airborne, on time, with correct weapons, across a 30-day campaign.",
         size=15, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(sl, 3, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Our Solution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "OUR SOLUTION", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "AirBase Ops — AI-Powered Digital Twin", size=36, bold=True, color=WHITE)

add_text(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
         "We digitized Saab's physical logistics board game, then augmented it with AI that transforms training into decision support.",
         size=17, color=LIGHT)

# Architecture boxes
layers = [
    ("FRONTEND", "React 19 + Tailwind CSS + Leaflet Maps", ACCENT, Inches(0.8)),
    ("BACKEND", "FastAPI + Python 3.11 + Pydantic v2", ACCENT2, Inches(4.8)),
    ("AI ENGINE", "OpenRouter LLM + Deterministic Fallback", GREEN, Inches(8.8)),
]

for label, tech, color, x in layers:
    add_rect(sl, x, Inches(2.5), Inches(3.6), Inches(1.8), BG_CARD)
    add_rect(sl, x, Inches(2.5), Inches(3.6), Inches(0.06), color)
    add_text(sl, x + Inches(0.3), Inches(2.7), Inches(3.0), Inches(0.4),
             label, size=13, bold=True, color=color)
    add_text(sl, x + Inches(0.3), Inches(3.1), Inches(3.0), Inches(1.0),
             tech, size=14, color=LIGHT)

# Feature modules under frontend
fe_items = [
    "3-column operations layout",
    "Interactive tactical map",
    "Real-time KPI dashboard",
    "AI recommendation cards",
]
for i, item in enumerate(fe_items):
    add_text(sl, Inches(1.1), Inches(3.45 + i * 0.28), Inches(3.0), Inches(0.3),
             f"  {item}", size=11, color=MUTED)

# Feature modules under backend
be_items = [
    "Turn engine (hourly sim)",
    "Resource manager (7 types)",
    "Maintenance system (dice)",
    "30-day campaign scenarios",
]
for i, item in enumerate(be_items):
    add_text(sl, Inches(5.1), Inches(3.45 + i * 0.28), Inches(3.0), Inches(0.3),
             f"  {item}", size=11, color=MUTED)

# Feature modules under AI
ai_items = [
    "Structured recommendations",
    "Mission allocation optimizer",
    "Conversational advisor",
    "Offline deterministic planner",
]
for i, item in enumerate(ai_items):
    add_text(sl, Inches(9.1), Inches(3.45 + i * 0.28), Inches(3.0), Inches(0.3),
             f"  {item}", size=11, color=MUTED)

# Key value props
props = [
    ("10 Gripen E", "Full fleet simulation with\nlifecycle state machine", ACCENT),
    ("30-Day Campaign", "Peace > Crisis > War\nwith dynamic events", ORANGE),
    ("4 Real-Time KPIs", "Readiness, throughput,\nturnaround, risk score", GREEN),
    ("AI Decision Support", "Explainable recommendations\nwith confidence scores", ACCENT2),
]

for i, (title, desc, color) in enumerate(props):
    x = Inches(0.8) + Inches(i * 3.15)
    y = Inches(5.0)
    add_rect(sl, x, y, Inches(2.8), Inches(1.8), BG_CARD)
    add_rect(sl, x, y, Inches(0.06), Inches(1.8), color)
    add_text(sl, x + Inches(0.25), y + Inches(0.15), Inches(2.4), Inches(0.4),
             title, size=16, bold=True, color=color)
    add_text(sl, x + Inches(0.25), y + Inches(0.6), Inches(2.4), Inches(1.0),
             desc, size=13, color=LIGHT)

slide_number(sl, 4, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Fleet Operations & Tactical Map
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "FEATURE SPOTLIGHT", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Fleet Operations & Tactical Map", size=36, bold=True, color=WHITE)

# Left — Fleet Board
add_rect(sl, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8), BG_CARD)
add_rect(sl, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.06), ACCENT)
add_text(sl, Inches(1.1), Inches(2.2), Inches(5.0), Inches(0.4),
         "FLEET OPERATIONS BOARD", size=13, bold=True, color=ACCENT)

fleet_features = [
    ("Real-time fleet overview", "All 10 aircraft visible at a glance with color-coded status"),
    ("Aircraft detail cards", "Click any aircraft for fuel, weapons, service hours, maintenance info"),
    ("Drag-and-assign", "Assign aircraft to missions directly from the fleet view"),
    ("Status indicators", "Green = ready, Blue = flying, Yellow = prepping, Red = maintenance"),
    ("Fuel level bars", "Visual fuel gauges updated every turn"),
    ("Service hour tracking", "Remaining hours until scheduled maintenance"),
]

for i, (title, desc) in enumerate(fleet_features):
    y = Inches(2.7 + i * 0.65)
    add_text(sl, Inches(1.3), y, Inches(5.0), Inches(0.3),
             title, size=14, bold=True, color=WHITE)
    add_text(sl, Inches(1.3), y + Inches(0.28), Inches(5.0), Inches(0.3),
             desc, size=11, color=MUTED)

# Right — Tactical Map
add_rect(sl, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.8), BG_CARD)
add_rect(sl, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.06), GREEN)
add_text(sl, Inches(7.1), Inches(2.2), Inches(5.0), Inches(0.4),
         "INTERACTIVE TACTICAL MAP", size=13, bold=True, color=GREEN)

map_features = [
    ("Leaflet.js interactive map", "Pan, zoom, and inspect the operational theater"),
    ("Base location indicator", "F 17 Kallinge shown with breathing glow effect"),
    ("Mission zones", "Color-coded circles: QRA (red), DCA (orange), RECCE (yellow)"),
    ("Aircraft positions", "Live markers showing aircraft locations and flight tracks"),
    ("Click-to-inspect", "Click any mission zone or aircraft for detailed information"),
    ("Auto-zoom", "Camera automatically fits all active missions and routes"),
]

for i, (title, desc) in enumerate(map_features):
    y = Inches(2.7 + i * 0.65)
    add_text(sl, Inches(7.3), y, Inches(5.0), Inches(0.3),
             title, size=14, bold=True, color=WHITE)
    add_text(sl, Inches(7.3), y + Inches(0.28), Inches(5.0), Inches(0.3),
             desc, size=11, color=MUTED)

slide_number(sl, 5, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — AI Decision Support
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "CORE INNOVATION", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "AI That Explains Its Reasoning", size=36, bold=True, color=WHITE)

add_text(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
         "Three layers of AI intelligence — from structured recommendations to free-form conversation.",
         size=17, color=LIGHT)

# Three AI layers
ai_layers = [
    ("1", "Recommendation Cards", ACCENT,
     "Structured, actionable suggestions",
     ["Specific aircraft-to-mission assignments",
      "Bullet-point reasoning for each action",
      "Expected effects on 4 KPIs with deltas",
      "Confidence score (0-100%)",
      "Stated assumptions & tradeoffs",
      "One-click \"Apply All\" execution"]),
    ("2", "Compare Mode", ACCENT2,
     "Non-destructive 6-hour simulation",
     ["Side-by-side: Your Plan vs AI-Optimized",
      "Metrics: missions, readiness, fuel, risk",
      "Green/red deltas for each metric",
      "Learn through comparison, not failure",
      "Zero risk — game state unchanged",
      "Builds trust in AI recommendations"]),
    ("3", "Chat Advisor", GREEN,
     "Natural language Q&A",
     ["\"Assess readiness\" — instant fleet briefing",
      "\"Fuel forecast\" — burn rate projection",
      "\"Maint priorities\" — repair sequencing",
      "Understands Swedish military terminology",
      "References specific aircraft & mission IDs",
      "Context-aware across campaign phases"]),
]

for i, (num, title, color, subtitle, bullets) in enumerate(ai_layers):
    x = Inches(0.8) + Inches(i * 4.0)
    # Card
    add_rect(sl, x, Inches(2.5), Inches(3.6), Inches(4.3), BG_CARD)
    add_rect(sl, x, Inches(2.5), Inches(3.6), Inches(0.06), color)
    # Number badge
    add_rect(sl, x + Inches(0.2), Inches(2.7), Inches(0.45), Inches(0.45), color)
    add_text(sl, x + Inches(0.2), Inches(2.72), Inches(0.45), Inches(0.45),
             num, size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    # Title
    add_text(sl, x + Inches(0.8), Inches(2.72), Inches(2.6), Inches(0.4),
             title, size=18, bold=True, color=color)
    # Subtitle
    add_text(sl, x + Inches(0.3), Inches(3.25), Inches(3.0), Inches(0.4),
             subtitle, size=13, color=MUTED)
    # Bullets
    for j, bullet in enumerate(bullets):
        add_text(sl, x + Inches(0.3), Inches(3.75 + j * 0.45), Inches(3.0), Inches(0.4),
                 f"  {bullet}", size=12, color=LIGHT)

slide_number(sl, 6, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Decision Impact Panel (KPIs)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "METRICS-DRIVEN DECISIONS", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "4 Real-Time KPIs — Always Visible", size=36, bold=True, color=WHITE)

add_text(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
         "Every decision shows immediate impact. Green arrows for improvements, red for degradation.",
         size=17, color=LIGHT)

kpis = [
    ("Fleet Readiness", "78%", GREEN,
     "Mission-capable + flying aircraft\nas percentage of total fleet",
     "(CAPABLE + ON_MISSION) / TOTAL"),
    ("Mission Throughput", "5", ACCENT,
     "Missions completable in the\nnext 6 hours given current fleet",
     "Available + returning + prepping"),
    ("Turnaround Delay", "142 min", ORANGE,
     "Average minutes until non-ready\naircraft become mission-capable",
     "Weighted avg of remaining hours"),
    ("Risk Score", "28", RED,
     "Composite index of fuel, UE,\nmaintenance burden, mission gaps",
     "4 components x 25 pts each"),
]

for i, (title, value, color, desc, formula) in enumerate(kpis):
    x = Inches(0.5) + Inches(i * 3.15)
    # Card
    add_rect(sl, x, Inches(2.5), Inches(2.9), Inches(4.2), BG_CARD)
    add_rect(sl, x, Inches(2.5), Inches(2.9), Inches(0.06), color)
    # Title
    add_text(sl, x + Inches(0.25), Inches(2.7), Inches(2.4), Inches(0.4),
             title, size=15, bold=True, color=color)
    # Value
    add_text(sl, x + Inches(0.25), Inches(3.2), Inches(2.4), Inches(0.8),
             value, size=48, bold=True, color=WHITE)
    # Delta example
    add_text(sl, x + Inches(0.25), Inches(4.15), Inches(2.4), Inches(0.4),
             "  +12% from last turn", size=12, color=GREEN)
    # Description
    add_text(sl, x + Inches(0.25), Inches(4.7), Inches(2.4), Inches(1.0),
             desc, size=13, color=LIGHT)
    # Formula
    add_rect(sl, x + Inches(0.15), Inches(5.7), Inches(2.6), Inches(0.6), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, x + Inches(0.25), Inches(5.75), Inches(2.4), Inches(0.5),
             formula, size=10, color=MUTED)

slide_number(sl, 7, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Game Engine (Faithful Digitization)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "FAITHFUL DIGITIZATION", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Every Board Game Mechanic — Digitized", size=36, bold=True, color=WHITE)

# Aircraft State Machine
add_rect(sl, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8), BG_CARD)
add_rect(sl, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.06), ACCENT2)
add_text(sl, Inches(1.1), Inches(2.2), Inches(5.2), Inches(0.4),
         "AIRCRAFT LIFECYCLE STATE MACHINE", size=13, bold=True, color=ACCENT2)

states = [
    ("HANGAR", "Not prepared", "gray", MUTED),
    ("PRE_FLIGHT", "4h prep, 1/3 fault chance", "yellow", YELLOW),
    ("MISSION_CAPABLE", "Ready to fly", "green", GREEN),
    ("ON_MISSION", "Currently flying", "blue", ACCENT),
    ("POST_FLIGHT", "6h, 50% maintenance", "orange", ORANGE),
    ("MAINTENANCE", "2-16h repair", "red", RED),
]

for i, (state, desc, _, color) in enumerate(states):
    y = Inches(2.8 + i * 0.6)
    # Color dot
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.05), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # State name
    add_text(sl, Inches(1.6), y, Inches(2.2), Inches(0.3),
             state, size=13, bold=True, color=WHITE, font_name="Consolas")
    # Description
    add_text(sl, Inches(3.9), y, Inches(2.5), Inches(0.3),
             desc, size=12, color=MUTED)

# Arrow flow text
add_text(sl, Inches(1.1), Inches(6.4), Inches(5.2), Inches(0.3),
         "HANGAR > PRE_FLIGHT > CAPABLE > MISSION > POST_FLIGHT > HANGAR",
         size=10, color=ACCENT2, font_name="Consolas", align=PP_ALIGN.CENTER)

# Resources panel
add_rect(sl, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4.8), BG_CARD)
add_rect(sl, Inches(7.0), Inches(2.0), Inches(5.4), Inches(0.06), ORANGE)
add_text(sl, Inches(7.3), Inches(2.2), Inches(5.0), Inches(0.4),
         "COMPLETE RESOURCE SIMULATION", size=13, bold=True, color=ORANGE)

resources = [
    ("Fuel Storage", "180,000L", "200L/flight-hr + 50L takeoff"),
    ("Missiles", "180 units", "10-100% lost per mission (dice)"),
    ("Bombs", "120 units", "Mission-type dependent consumption"),
    ("Recon Pods", "10 units", "Required for RECCE missions"),
    ("Spare Parts", "60 units", "Consumed by maintenance repairs"),
    ("Exchange Units", "16 units", "30-day MRO repair cycle"),
    ("Maint. Crews", "6 (3 on duty)", "8-hour rotating shifts"),
]

for i, (name, amount, consumption) in enumerate(resources):
    y = Inches(2.8 + i * 0.55)
    add_text(sl, Inches(7.3), y, Inches(2.0), Inches(0.3),
             name, size=13, bold=True, color=WHITE)
    add_text(sl, Inches(9.3), y, Inches(1.2), Inches(0.3),
             amount, size=13, color=ACCENT2)
    add_text(sl, Inches(10.3), y, Inches(2.0), Inches(0.3),
             consumption, size=10, color=MUTED)

# Campaign phases at bottom
add_text(sl, Inches(7.3), Inches(6.0), Inches(5.0), Inches(0.3),
         "CAMPAIGN PHASES", size=11, bold=True, color=ORANGE)

phases = [
    ("Days 1-10: PEACE", GREEN, Inches(7.3)),
    ("Days 11-20: CRISIS", YELLOW, Inches(9.3)),
    ("Days 21-30: WAR", RED, Inches(11.0)),
]
for text, color, x in phases:
    add_rect(sl, x, Inches(6.35), Inches(1.5), Inches(0.35), color)
    add_text(sl, x + Inches(0.05), Inches(6.37), Inches(1.4), Inches(0.3),
             text, size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

slide_number(sl, 8, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — Before vs After
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "TRANSFORMATION", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Board Game vs AirBase Ops", size=36, bold=True, color=WHITE)

# Comparison table
headers = [("Dimension", Inches(0.8), Inches(2.8)),
           ("Physical Board Game", Inches(4.5), Inches(3.2)),
           ("AirBase Ops", Inches(8.8), Inches(3.2))]

# Column backgrounds
add_rect(sl, Inches(0.8), Inches(2.2), Inches(3.4), Inches(0.5), BG_CARD)
add_rect(sl, Inches(4.3), Inches(2.2), Inches(4.2), Inches(0.5), RGBColor(0x3B, 0x18, 0x18))
add_rect(sl, Inches(8.6), Inches(2.2), Inches(4.0), Inches(0.5), RGBColor(0x0C, 0x2D, 0x1E))

add_text(sl, Inches(1.0), Inches(2.25), Inches(3.0), Inches(0.4),
         "DIMENSION", size=12, bold=True, color=MUTED)
add_text(sl, Inches(4.5), Inches(2.25), Inches(3.8), Inches(0.4),
         "PHYSICAL BOARD GAME", size=12, bold=True, color=RED)
add_text(sl, Inches(8.8), Inches(2.25), Inches(3.6), Inches(0.4),
         "AIRBASE OPS", size=12, bold=True, color=GREEN)

comparisons = [
    ("Players needed", "8-15 people", "1+ (solo or team)"),
    ("Session setup", "Hours of preparation", "Click and play in seconds"),
    ("Game speed", "4-8 hours per day", "Real-time, adjustable speed"),
    ("Decision support", "Human mentors only", "AI advisor + what-if sim"),
    ("Data capture", "None", "Full decision audit trail"),
    ("Availability", "Scheduled events only", "24/7, any browser"),
    ("Scalability", "1 session at a time", "Unlimited concurrent"),
    ("What-if analysis", "Impossible", "Compare Mode with 6h forecast"),
]

for i, (dim, old, new) in enumerate(comparisons):
    y = Inches(2.85 + i * 0.52)
    # Alternating row bg
    if i % 2 == 0:
        add_rect(sl, Inches(0.8), y - Inches(0.05), Inches(11.8), Inches(0.48), BG_CARD)
    add_text(sl, Inches(1.0), y, Inches(3.0), Inches(0.35),
             dim, size=14, bold=True, color=WHITE)
    add_text(sl, Inches(4.5), y, Inches(3.8), Inches(0.35),
             old, size=14, color=RED)
    add_text(sl, Inches(8.8), y, Inches(3.6), Inches(0.35),
             new, size=14, color=GREEN)

slide_number(sl, 9, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — Business Case
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "BUSINESS CASE", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Why Invest in This Product?", size=36, bold=True, color=WHITE)

# Five value props
biz_items = [
    ("Training Scalability", ACCENT,
     "Physical game: 1 session, 1 room, 8+ people\nAirBase Ops: Unlimited sessions, any browser, solo or team\nTraining frequency goes from 2x/year to continuous"),
    ("Customer Value Amplifier", GREEN,
     "Every Gripen export customer needs base ops training\nPackage with aircraft sales as training differentiator\nStrengthens Saab's total support offering"),
    ("Data-Driven Insights", ACCENT2,
     "Every decision logged with timestamps and state\nAfter-action review and performance analytics\nCompare cohort performance over time"),
    ("AI Augmentation Path", ORANGE,
     "Today: Training tool with AI suggestions\nTomorrow: Operational decision support for real bases\nSame engine, real fleet data instead of simulation"),
    ("Edge Deployment Ready", YELLOW,
     "Python + React = lightweight, standard stack\nDeterministic fallback = works offline\nNo specialized hardware — runs on field laptops"),
]

for i, (title, color, desc) in enumerate(biz_items):
    if i < 3:
        x = Inches(0.8) + Inches(i * 4.0)
        y = Inches(1.9)
    else:
        x = Inches(0.8) + Inches((i - 3) * 4.0) + Inches(2.0)
        y = Inches(4.5)

    add_rect(sl, x, y, Inches(3.6), Inches(2.2), BG_CARD)
    add_rect(sl, x, y, Inches(3.6), Inches(0.06), color)
    add_text(sl, x + Inches(0.3), y + Inches(0.2), Inches(3.0), Inches(0.4),
             title, size=16, bold=True, color=color)
    add_text(sl, x + Inches(0.3), y + Inches(0.65), Inches(3.0), Inches(1.4),
             desc, size=12, color=LIGHT)

slide_number(sl, 10, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — SaaS Revenue Model
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "REVENUE MODEL", size=14, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "SaaS Tiers & Future Roadmap", size=36, bold=True, color=WHITE)

# Three pricing tiers
tiers = [
    ("Training", "Core simulation, scenario library,\nsingle-player mode, basic analytics",
     "Per-seat license", ACCENT, [
         "10 Gripen fleet simulation",
         "30-day campaign scenarios",
         "Aircraft lifecycle management",
         "Resource tracking dashboard",
     ]),
    ("Advanced", "AI advisor, compare mode,\nrecommendation cards, full analytics",
     "Per-base subscription", ACCENT2, [
         "Everything in Training, plus:",
         "AI-powered recommendations",
         "Compare Mode (what-if)",
         "Decision audit trail & export",
     ]),
    ("Enterprise", "Custom scenarios, API integration,\nmulti-base networking, multiplayer",
     "Enterprise agreement", GREEN, [
         "Everything in Advanced, plus:",
         "Custom scenario editor",
         "Multi-base coordination",
         "Multiplayer roles (AOC/BC/KlT)",
     ]),
]

for i, (name, desc, price, color, features) in enumerate(tiers):
    x = Inches(0.8) + Inches(i * 4.0)
    # Card
    add_rect(sl, x, Inches(2.0), Inches(3.6), Inches(4.5), BG_CARD)
    add_rect(sl, x, Inches(2.0), Inches(3.6), Inches(0.06), color)
    # Tier name
    add_text(sl, x + Inches(0.3), Inches(2.25), Inches(3.0), Inches(0.5),
             name, size=24, bold=True, color=color)
    # Description
    add_text(sl, x + Inches(0.3), Inches(2.85), Inches(3.0), Inches(0.8),
             desc, size=12, color=MUTED)
    # Price model
    add_rect(sl, x + Inches(0.2), Inches(3.75), Inches(3.2), Inches(0.4), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, x + Inches(0.3), Inches(3.78), Inches(3.0), Inches(0.35),
             price, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Feature list
    for j, feat in enumerate(features):
        add_text(sl, x + Inches(0.3), Inches(4.35 + j * 0.4), Inches(3.0), Inches(0.35),
                 f"  {feat}", size=12, color=LIGHT)

# Future roadmap bar
add_rect(sl, Inches(0.6), H - Inches(0.9), Inches(12), Inches(0.6), BG_CARD)
add_text(sl, Inches(0.9), H - Inches(0.85), Inches(11.4), Inches(0.15),
         "FUTURE ROADMAP", size=10, bold=True, color=ACCENT)
add_text(sl, Inches(0.9), H - Inches(0.6), Inches(11.4), Inches(0.3),
         "Mobile/tablet support  |  VR/AR immersive training  |  Real fleet data integration  |  NATO-standard interoperability  |  Predictive maintenance AI",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(sl, 11, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — Closing
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# Big statement
add_text(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
         "We didn't just digitize a board game.", size=32, color=LIGHT,
         align=PP_ALIGN.CENTER)

add_text(sl, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
         "We built the foundation for Saab's\nnext-generation air base command system.",
         size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three closing pillars
pillars = [
    ("Proven Training\nMethodology", "Faithful digitization of\nSaab's board game", ACCENT),
    ("AI-Augmented\nDecision Support", "Explainable AI that\nbuilds trust & skill", ACCENT2),
    ("Production-Ready\nArchitecture", "Edge-deployable, offline,\nscalable to enterprise", GREEN),
]

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(1.5) + Inches(i * 3.8)
    add_rect(sl, x, Inches(4.5), Inches(3.2), Inches(1.6), BG_CARD)
    add_rect(sl, x, Inches(4.5), Inches(3.2), Inches(0.06), color)
    add_text(sl, x + Inches(0.25), Inches(4.65), Inches(2.7), Inches(0.7),
             title, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.25), Inches(5.35), Inches(2.7), Inches(0.6),
             desc, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# CTA
add_text(sl, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "AIRBASE OPS", size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(sl, Inches(1.5), Inches(7.0), Inches(10), Inches(0.3),
         "Saab Smart Air Base Hackathon  |  March 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(sl, 12, TOTAL_SLIDES)


# ── Save ─────────────────────────────────────────────
output_path = "/Users/omarali/Web/SaabHackathon/AirBaseOps_Pitch.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
